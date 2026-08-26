"""Deck Agent - presentation generation for marketing content.

Supports two output formats:
- Google Slides (default) via direct Google Slides API
- PowerPoint (.pptx) via python-pptx for offline/download use
"""

from __future__ import annotations

import json
import re
from io import BytesIO
from pathlib import Path
from typing import Any, Dict, List, Optional
from uuid import uuid4

import structlog
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from ..db.database import Database
from ..db.repositories import DraftRepo
from ..llm.base import BaseLLM, Message
from ..quality import QualityCriterion, RefinementConfig, RefinementLoop
from ..quality.hooks import build_hook_chain, make_em_dash_hook
from ..workspace.manager import WorkspaceManager
from .base import BaseAgent

logger = structlog.get_logger()

# Quality criteria for deck content structure
_DECK_CRITERIA = [
    QualityCriterion(
        name="slide_conciseness",
        description="Slides are concise with 6-8 words per bullet, max 5 bullets per slide",
    ),
    QualityCriterion(
        name="visual_hierarchy",
        description="Clear visual hierarchy with varied slide types for engagement",
    ),
    QualityCriterion(
        name="narrative_flow",
        description="Logical narrative arc: hook, problem, solution, evidence, CTA",
    ),
    QualityCriterion(
        name="brand_consistency",
        description="Content aligns with brand voice, tone, and messaging guidelines",
    ),
    QualityCriterion(
        name="data_clarity",
        description="Data, stats, and claims are clearly presented and easy to understand",
    ),
]

# ── Default output directory ──────────────────────────────────────────────
DECKS_DIR = Path("data/decks")

# ── Brand color palettes ─────────────────────────────────────────────────
# Used for .pptx generation (RGBColor objects) and Google Slides (converted to floats)
DEFAULT_PALETTE = {
    "primary": RGBColor(0x1A, 0x73, 0xE8),  # Blue
    "secondary": RGBColor(0x34, 0xA8, 0x53),  # Green
    "accent": RGBColor(0xFB, 0xBC, 0x04),  # Amber
    "dark": RGBColor(0x20, 0x2A, 0x36),  # Dark navy
    "light": RGBColor(0xF8, 0xF9, 0xFA),  # Off-white
    "text": RGBColor(0x33, 0x33, 0x33),  # Dark gray
    "text_light": RGBColor(0xFF, 0xFF, 0xFF),  # White
    "muted": RGBColor(0x6C, 0x75, 0x7D),  # Gray
}

BRAND_PALETTES: Dict[str, Dict[str, RGBColor]] = {
    "uhg": {
        "primary": RGBColor(0x00, 0x2F, 0x6C),  # UnitedHealth deep blue
        "secondary": RGBColor(0xF2, 0xA9, 0x00),  # UHG gold/orange accent
        "accent": RGBColor(0x00, 0x82, 0xC8),  # Optum blue accent
        "dark": RGBColor(0x0A, 0x1F, 0x44),  # Dark navy
        "light": RGBColor(0xE8, 0xF1, 0xFA),  # Light blue
        "text": RGBColor(0x21, 0x21, 0x21),
        "text_light": RGBColor(0xFF, 0xFF, 0xFF),
        "muted": RGBColor(0x6C, 0x75, 0x7D),
    },
}

# ── Slide layout constants ────────────────────────────────────────────────
SLIDE_WIDTH = Inches(13.333)  # 16:9 widescreen
SLIDE_HEIGHT = Inches(7.5)

# Google Slides page dimensions in EMU (9144000 x 5143500 for 10" x 5.625")
GS_PAGE_WIDTH = 9144000
GS_PAGE_HEIGHT = 5143500

# Google Slides layout constants (inches) — 10" x 5.625" canvas
_GS_W = 10.0  # page width
_GS_H = 5.625  # page height
_GS_MX = 0.6  # horizontal margin
_GS_MY = 0.4  # vertical margin

# White background for all Google Slides (float RGB)
_GS_WHITE = {"red": 1.0, "green": 1.0, "blue": 1.0}


class DeckAgent(BaseAgent):
    """Creates presentations for marketing campaigns.

    Produces polished slide decks with:
    - LLM-generated content structure and copy
    - Brand-appropriate color schemes and styling
    - Multiple slide types (title, content, two-column, stats, image, CTA)
    - Google Slides as primary output (direct API)
    - PowerPoint (.pptx) as secondary/offline format
    """

    agent_id = "deck_agent"
    agent_name = "Deck Agent"

    def __init__(
        self,
        llm: BaseLLM,
        db: Database,
        workspace_manager: WorkspaceManager,
        google_credentials_path: str = "",
        google_oauth_token_path: str = "",
        decks_dir: Optional[Path] = None,
        media_storage: Optional[Any] = None,
        proofreader: Optional[Any] = None,
        scanning_llm: Optional[BaseLLM] = None,
        max_iterations: int = 10,
    ) -> None:
        self._drafts = DraftRepo(db)
        self._workspace_mgr = workspace_manager
        self._google_credentials_path = google_credentials_path
        self._google_oauth_token_path = google_oauth_token_path
        self._decks_dir = decks_dir or DECKS_DIR
        self._decks_dir.mkdir(parents=True, exist_ok=True)
        self._slides_service = None
        self._drive_service = None
        self._media_storage = media_storage
        self._proofreader = proofreader
        self._scanning_llm = scanning_llm
        super().__init__(llm=llm, db=db, max_iterations=max_iterations)

    def _get_google_services(self) -> bool:
        """Lazily initialize Google Slides and Drive API services.

        Returns True if services are available, False otherwise.
        Uses OAuth2 token (preferred) or service account credentials (fallback).
        """
        if self._slides_service is not None:
            return True

        try:
            from googleapiclient.discovery import build

            from ..google_auth import get_google_credentials

            SCOPES = [
                "https://www.googleapis.com/auth/presentations",
                "https://www.googleapis.com/auth/drive.file",
            ]
            # Force OAuth — service account lacks Presentations API permission.
            credentials = get_google_credentials(
                oauth_token_path=self._google_oauth_token_path,
                service_account_path="",
                scopes=SCOPES,
            )
            if credentials is None:
                return False

            self._slides_service = build("slides", "v1", credentials=credentials)
            self._drive_service = build("drive", "v3", credentials=credentials)
            logger.info("google_services_initialized")
            return True
        except Exception as e:
            logger.error("google_services_init_failed", error=str(e))
            return False

    def get_system_prompt(self, workspace_id: Optional[str] = None) -> str:
        ws_label = workspace_id or "all workspaces"
        gslides_available = bool(self._google_credentials_path)
        return f"""You are the Deck Agent for {ws_label}.
You create professional presentations for marketing campaigns, pitches,
reports, and proposals.

Output formats:
- Google Slides (preferred) — credentials: {"configured" if gslides_available else "not configured"}
- PowerPoint (.pptx) for offline/download use
- Default is "google_slides" when credentials are configured, otherwise "pptx"

CRITICAL RULES:
- You MUST ALWAYS call the create_deck, create_deck_from_content, or create_pitch_deck tool. NEVER generate slide content as text in your response.
- Your ONLY job is to call a tool that creates a real presentation and return the link. Do NOT describe slides in your response.
- If the tool returns "not_configured", report that Google credentials need to be set up. Do NOT fall back to generating content as text.
- After creating a deck, your response should be brief: state that the presentation was created and include the Google Slides URL.
- ABBREVIATIONS: Always spell out abbreviations/acronyms on first use with the abbreviation in parentheses — e.g., "Total Addressable Market (TAM)". After the first occurrence, use the abbreviation only.
- ANTI-AI WRITING (MANDATORY): NEVER use em dashes (—) or en dashes (–). NEVER bold words mid-sentence. Use commas, periods, semicolons, or " - " instead. Bolding is only for headers or standalone labels.

You can:
- Generate complete slide decks from a topic or brief
- Create individual slides of specific types
- Convert existing content (blogs, newsletters) into presentations
- Use brand colors and styling per workspace

When creating decks:
- Keep slides concise (6-8 words per bullet, max 5 bullets per slide)
- Use a mix of slide types for visual variety
- Include speaker notes for each slide
- Structure with a clear narrative arc: hook, problem, solution, evidence, CTA
- Default to 8-12 slides unless specified otherwise

Available slide types:
- title: Title slide with subtitle
- section: Section divider slide
- content: Bullet points with optional header
- two_column: Side-by-side content comparison
- stats: Key metrics / statistics highlight
- image: Full or half-image slide with caption
- quote: Large quote with attribution
- cta: Call-to-action closing slide
Current workspace: {ws_label}"""

    def register_tools(self) -> None:

        @self.tool_registry.register(
            name="create_deck",
            description=(
                "Generate a complete presentation from a topic or brief. "
                "Defaults to Google Slides when configured, otherwise .pptx. "
                "Set output_format to 'google_slides', 'pptx', or 'both'."
            ),
        )
        async def create_deck(
            topic: str,
            workspace_id: str,
            purpose: str = "marketing",
            audience: str = "general",
            num_slides: int = 10,
            include_speaker_notes: bool = True,
            output_format: str = "",
        ) -> Dict[str, Any]:
            fmt = self._resolve_format(output_format)
            brand_voice = await self._workspace_mgr.get_brand_voice(workspace_id)
            palette = _get_palette(workspace_id)

            structure = await self._generate_deck_structure(
                topic=topic,
                purpose=purpose,
                audience=audience,
                num_slides=num_slides,
                brand_voice=brand_voice,
                include_speaker_notes=include_speaker_notes,
            )

            if self._proofreader:
                try:
                    structure, _ = await self._proofreader.proofread_dict_fields(
                        structure,
                        fields=[
                            "title",
                            "subtitle",
                            "bullets",
                            "content",
                            "quote",
                            "attribution",
                            "left_column",
                            "right_column",
                            "left_title",
                            "right_title",
                            "caption",
                            "speaker_notes",
                            "cta",
                            "details",
                        ],
                    )
                except Exception as e:
                    logger.warning("deck_proofread_failed", error=str(e))

            result = await self._build_and_export(
                structure=structure,
                workspace_id=workspace_id,
                palette=palette,
                output_format=fmt,
            )

            return {
                **result,
                "topic": topic,
                "num_slides": len(structure.get("slides", [])),
                "purpose": purpose,
                "workspace_id": workspace_id,
                "output_format": fmt,
            }

        @self.tool_registry.register(
            name="create_deck_from_content",
            description=(
                "Generate a presentation from existing content (e.g., a blog post, "
                "newsletter, or research summary). Converts the content into slides."
            ),
        )
        async def create_deck_from_content(
            content: str,
            workspace_id: str,
            title: str = "",
            purpose: str = "presentation",
            audience: str = "general",
            output_format: str = "",
        ) -> Dict[str, Any]:
            fmt = self._resolve_format(output_format)
            brand_voice = await self._workspace_mgr.get_brand_voice(workspace_id)
            palette = _get_palette(workspace_id)

            prompt = f"""Convert the following content into a slide deck structure.

CONTENT:
{content[:3000]}

TITLE: {title or "(generate an appropriate title)"}
PURPOSE: {purpose}
AUDIENCE: {audience}

BRAND CONTEXT:
{brand_voice[:300] if brand_voice else "No brand voice specified."}

Create a well-structured slide deck (8-12 slides) that presents this content
in a compelling visual format. Use varied slide types.

{_structure_format_instructions()}"""

            response = await self.llm.complete(
                messages=[Message(role="user", content=prompt)],
                temperature=0.7,
            )

            structure = _parse_structure(response.content)

            if self._proofreader:
                try:
                    structure, _ = await self._proofreader.proofread_dict_fields(
                        structure,
                        fields=[
                            "title",
                            "subtitle",
                            "bullets",
                            "content",
                            "quote",
                            "attribution",
                            "left_column",
                            "right_column",
                            "left_title",
                            "right_title",
                            "caption",
                            "speaker_notes",
                            "cta",
                            "details",
                        ],
                    )
                except Exception as e:
                    logger.warning("deck_proofread_failed", error=str(e))

            result = await self._build_and_export(
                structure=structure,
                workspace_id=workspace_id,
                palette=palette,
                output_format=fmt,
            )

            return {
                **result,
                "title": structure.get("title", title),
                "num_slides": len(structure.get("slides", [])),
                "purpose": purpose,
                "workspace_id": workspace_id,
                "output_format": fmt,
            }

        @self.tool_registry.register(
            name="create_pitch_deck",
            description=(
                "Generate a pitch deck for a product, service, or campaign. "
                "Follows a standard pitch structure: problem, solution, market, "
                "traction, team, ask."
            ),
        )
        async def create_pitch_deck(
            company_or_product: str,
            workspace_id: str,
            value_prop: str = "",
            target_market: str = "",
            key_metrics: str = "",
            output_format: str = "",
        ) -> Dict[str, Any]:
            fmt = self._resolve_format(output_format)
            brand_voice = await self._workspace_mgr.get_brand_voice(workspace_id)
            palette = _get_palette(workspace_id)

            prompt = f"""Create a pitch deck structure for:

COMPANY/PRODUCT: {company_or_product}
VALUE PROPOSITION: {value_prop or "(determine from context)"}
TARGET MARKET: {target_market or "(determine from context)"}
KEY METRICS: {key_metrics or "(use placeholders like [INSERT METRIC])"}

BRAND CONTEXT:
{brand_voice[:300] if brand_voice else "No brand voice specified."}

Follow this standard pitch structure (10-12 slides):
1. Title slide - company name, tagline
2. Problem - the pain point you solve
3. Solution - your product/approach
4. How it works - 3-step process or demo
5. Market opportunity - TAM/SAM/SOM
6. Traction - key metrics and milestones
7. Business model - how you make money
8. Competition - your differentiators
9. Team - key people
10. The ask - what you need
11. CTA / contact slide

{_structure_format_instructions()}"""

            response = await self.llm.complete(
                messages=[Message(role="user", content=prompt)],
                temperature=0.7,
            )

            structure = _parse_structure(response.content)

            if self._proofreader:
                try:
                    structure, _ = await self._proofreader.proofread_dict_fields(
                        structure,
                        fields=[
                            "title",
                            "subtitle",
                            "bullets",
                            "content",
                            "quote",
                            "attribution",
                            "left_column",
                            "right_column",
                            "left_title",
                            "right_title",
                            "caption",
                            "speaker_notes",
                            "cta",
                            "details",
                        ],
                    )
                except Exception as e:
                    logger.warning("deck_proofread_failed", error=str(e))

            result = await self._build_and_export(
                structure=structure,
                workspace_id=workspace_id,
                palette=palette,
                output_format=fmt,
            )

            return {
                **result,
                "title": structure.get("title", company_or_product),
                "num_slides": len(structure.get("slides", [])),
                "purpose": "pitch",
                "workspace_id": workspace_id,
                "output_format": fmt,
            }

        @self.tool_registry.register(
            name="list_decks",
            description="List all generated deck files for a workspace.",
        )
        async def list_decks(workspace_id: str = "") -> Dict[str, Any]:
            decks = []
            for f in sorted(self._decks_dir.glob("*.pptx")):
                if not workspace_id or workspace_id in f.stem:
                    decks.append(
                        {
                            "file": f.name,
                            "path": str(f),
                            "size_kb": round(f.stat().st_size / 1024, 1),
                        }
                    )
            return {"decks": decks, "count": len(decks)}

    # ── Internal methods ──────────────────────────────────────────────────

    async def _generate_deck_structure(
        self,
        topic: str,
        purpose: str,
        audience: str,
        num_slides: int,
        brand_voice: str,
        include_speaker_notes: bool,
    ) -> Dict[str, Any]:
        """Use the LLM to generate a structured JSON deck outline."""
        notes_instruction = ""
        if include_speaker_notes:
            notes_instruction = '- "speaker_notes": talking points for the presenter'

        prompt = f"""Create a {num_slides}-slide presentation structure.

TOPIC: {topic}
PURPOSE: {purpose}
AUDIENCE: {audience}

BRAND CONTEXT:
{brand_voice[:400] if brand_voice else "No brand voice specified."}

{_structure_format_instructions(notes_instruction)}"""

        config = RefinementConfig(
            criteria=_DECK_CRITERIA,
            max_iterations=2,
            quality_threshold=7.0,
            log_prefix="deck_structure",
        )
        hooks = build_hook_chain(make_em_dash_hook())
        loop = RefinementLoop(
            config=config,
            generation_llm=self.llm,
            critique_llm=self._scanning_llm,
        )
        result = await loop.generate_and_refine(prompt, brand_voice or "", hooks)

        return _parse_structure(result.content)

    def _resolve_format(self, requested: str) -> str:
        """Determine output format. Default to google_slides when credentials exist.

        Accepted values:
            ""               auto: google_slides if creds, otherwise pptx
            "google_slides"  Google Slides only (falls back to pptx on failure)
            "pptx"           Microsoft PowerPoint only
            "both"           Build both formats
        """
        if requested:
            r = requested.lower().strip()
            if r in ("google", "gslide", "gslides", "slides"):
                return "google_slides"
            if r in ("powerpoint", "office", "ms_powerpoint", "microsoft_powerpoint"):
                return "pptx"
            return r
        if self._google_credentials_path and Path(self._google_credentials_path).exists():
            return "google_slides"
        return "pptx"

    async def _build_and_export(
        self,
        structure: Dict[str, Any],
        workspace_id: str,
        palette: Dict[str, RGBColor],
        output_format: str,
    ) -> Dict[str, Any]:
        """Build the deck and export in the requested format(s).

        output_format values:
            "google_slides"  Google Slides only (falls back to .pptx on failure)
            "pptx"           Microsoft PowerPoint only
            "both"           Build both formats and return combined result
        """
        if output_format == "both":
            gslides_result = self._build_google_slides(structure, workspace_id, palette)
            pptx_result = await self._build_pptx_and_upload(structure, workspace_id, palette)
            merged: Dict[str, Any] = {"status": "created"}
            if gslides_result.get("status") == "created":
                if gslides_result.get("google_slides_url"):
                    merged["google_slides_url"] = gslides_result["google_slides_url"]
                if gslides_result.get("google_slides_id"):
                    merged["google_slides_id"] = gslides_result["google_slides_id"]
            if pptx_result.get("status") == "created":
                merged["file_path"] = pptx_result.get("file_path")
                if pptx_result.get("cdn_url"):
                    merged["cdn_url"] = pptx_result["cdn_url"]
                if pptx_result.get("asset_id"):
                    merged["asset_id"] = pptx_result["asset_id"]
            if (
                gslides_result.get("status") != "created"
                and pptx_result.get("status") != "created"
            ):
                merged["status"] = "error"
                merged["message"] = "Both Google Slides and .pptx generation failed"
            return merged

        if output_format == "google_slides":
            gslides_result = self._build_google_slides(structure, workspace_id, palette)
            if gslides_result.get("status") == "created":
                return gslides_result
            logger.warning(
                "google_slides_failed_falling_back_to_pptx",
                reason=gslides_result.get("message", "unknown"),
            )

        # Build .pptx locally (primary or fallback)
        return await self._build_pptx_and_upload(structure, workspace_id, palette)

    async def _build_pptx_and_upload(
        self,
        structure: Dict[str, Any],
        workspace_id: str,
        palette: Dict[str, RGBColor],
    ) -> Dict[str, Any]:
        """Build a .pptx file and upload to media storage if configured."""
        file_path = self._build_pptx(
            structure=structure,
            workspace_id=workspace_id,
            palette=palette,
        )

        result: Dict[str, Any] = {
            "file_path": str(file_path),
            "status": "created",
        }

        if self._media_storage:
            try:
                asset = await self._media_storage.upload_from_file(
                    file_path=str(file_path),
                    workspace_id=workspace_id,
                    media_type="deck",
                    source_agent=self.agent_id,
                    prompt=structure.get("title", ""),
                )
                result["cdn_url"] = asset.cdn_url
                result["asset_id"] = asset.id
            except Exception as e:
                logger.warning("deck_media_upload_failed", error=str(e))

        return result

    # ── Google Slides API builder ─────────────────────────────────────────

    def _build_google_slides(
        self,
        structure: Dict[str, Any],
        workspace_id: str,
        palette: Dict[str, RGBColor],
    ) -> Dict[str, Any]:
        """Create a Google Slides presentation via the Slides API directly."""
        if not self._get_google_services():
            return {
                "status": "not_configured",
                "message": (
                    "Google credentials not configured. "
                    "Set GOOGLE_CREDENTIALS_PATH in .env to a service account JSON file."
                ),
            }

        slides_data = structure.get("slides", [])
        if not slides_data:
            slides_data = [
                {
                    "type": "title",
                    "title": structure.get("title", "Untitled"),
                    "subtitle": structure.get("subtitle", ""),
                }
            ]

        title = structure.get("title", "Untitled Presentation")
        pal = _palette_to_floats(palette)

        try:
            # 1. Create an empty presentation
            presentation = (
                self._slides_service.presentations().create(body={"title": title}).execute()
            )
            presentation_id = presentation["presentationId"]

            # The new presentation has one blank slide by default — get its ID
            default_slide_id = presentation["slides"][0]["objectId"]

            # 1b. Delete placeholder shapes from the default slide
            #     (removes "Click to add title" / "Click to add subtitle" boxes)
            placeholder_reqs: List[Dict[str, Any]] = []
            for page_element in presentation["slides"][0].get("pageElements", []):
                if "shape" in page_element and "placeholder" in page_element["shape"]:
                    placeholder_reqs.append(
                        {"deleteObject": {"objectId": page_element["objectId"]}}
                    )
            if placeholder_reqs:
                self._slides_service.presentations().batchUpdate(
                    presentationId=presentation_id,
                    body={"requests": placeholder_reqs},
                ).execute()

            # 2. Build all batch update requests
            requests: List[Dict[str, Any]] = []

            for i, sd in enumerate(slides_data):
                slide_type = sd.get("type", "content")

                if i == 0:
                    # Reuse the default blank slide for the first slide
                    slide_id = default_slide_id
                else:
                    # Create a new blank slide
                    slide_id = f"slide_{uuid4().hex[:8]}"
                    requests.append(
                        {
                            "createSlide": {
                                "objectId": slide_id,
                                "insertionIndex": i,
                            }
                        }
                    )

                # Set white background on every slide
                requests.append(_gs_set_bg(slide_id, _GS_WHITE))

                # Build slide content requests based on type
                builder = _GS_SLIDE_BUILDERS.get(slide_type, _gs_build_content)
                slide_requests = builder(slide_id, sd, pal, i)
                requests.extend(slide_requests)

            # 3. Execute all requests in a single batch
            if requests:
                self._slides_service.presentations().batchUpdate(
                    presentationId=presentation_id,
                    body={"requests": requests},
                ).execute()

            # 4. Share with authorized users only (restricted access)
            from ..google_auth import share_file_restricted

            share_file_restricted(self._drive_service, presentation_id)

            # 5. Move into CMO Agent folder
            from ..google_auth import ensure_drive_folder, move_file_to_folder

            folder_id = ensure_drive_folder(self._drive_service)
            if folder_id:
                move_file_to_folder(self._drive_service, presentation_id, folder_id)

            slides_url = f"https://docs.google.com/presentation/d/{presentation_id}/edit"

            logger.info(
                "google_slides_created",
                presentation_id=presentation_id,
                slides=len(slides_data),
                workspace=workspace_id,
            )

            return {
                "status": "created",
                "google_slides_url": slides_url,
                "google_slides_id": presentation_id,
            }

        except Exception as e:
            logger.error("google_slides_build_failed", error=str(e))
            return {"status": "error", "message": str(e)}

    def _update_google_slides(
        self,
        presentation_id: str,
        structure: Dict[str, Any],
        workspace_id: str,
        palette: Dict[str, RGBColor],
    ) -> Dict[str, Any]:
        """Update an existing Google Slides presentation by clearing and rebuilding."""
        if not self._get_google_services():
            return {
                "status": "not_configured",
                "message": "Google credentials not configured.",
            }

        slides_data = structure.get("slides", [])
        if not slides_data:
            slides_data = [
                {
                    "type": "title",
                    "title": structure.get("title", "Untitled"),
                    "subtitle": structure.get("subtitle", ""),
                }
            ]

        pal = _palette_to_floats(palette)

        try:
            # 1. Get existing presentation
            pres = (
                self._slides_service.presentations().get(presentationId=presentation_id).execute()
            )
            existing_slides = pres.get("slides", [])

            # 2. Delete all slides except the first one
            delete_reqs: List[Dict[str, Any]] = []
            for slide in existing_slides[1:]:
                delete_reqs.append({"deleteObject": {"objectId": slide["objectId"]}})
            if delete_reqs:
                self._slides_service.presentations().batchUpdate(
                    presentationId=presentation_id,
                    body={"requests": delete_reqs},
                ).execute()

            # 3. Clear all elements from the first slide
            pres = (
                self._slides_service.presentations().get(presentationId=presentation_id).execute()
            )
            first_slide = pres["slides"][0]
            first_slide_id = first_slide["objectId"]
            clear_reqs: List[Dict[str, Any]] = []
            for elem in first_slide.get("pageElements", []):
                clear_reqs.append({"deleteObject": {"objectId": elem["objectId"]}})
            if clear_reqs:
                self._slides_service.presentations().batchUpdate(
                    presentationId=presentation_id,
                    body={"requests": clear_reqs},
                ).execute()

            # 4. Rebuild all slides
            requests: List[Dict[str, Any]] = []

            for i, sd in enumerate(slides_data):
                slide_type = sd.get("type", "content")

                if i == 0:
                    slide_id = first_slide_id
                else:
                    slide_id = f"slide_{uuid4().hex[:8]}"
                    requests.append(
                        {
                            "createSlide": {
                                "objectId": slide_id,
                                "insertionIndex": i,
                            }
                        }
                    )

                # Set white background on every slide
                requests.append(_gs_set_bg(slide_id, _GS_WHITE))

                builder = _GS_SLIDE_BUILDERS.get(slide_type, _gs_build_content)
                slide_requests = builder(slide_id, sd, pal, i)
                requests.extend(slide_requests)

            if requests:
                self._slides_service.presentations().batchUpdate(
                    presentationId=presentation_id,
                    body={"requests": requests},
                ).execute()

            slides_url = f"https://docs.google.com/presentation/d/{presentation_id}/edit"

            logger.info(
                "google_slides_updated",
                presentation_id=presentation_id,
                slides=len(slides_data),
                workspace=workspace_id,
            )

            return {
                "status": "updated",
                "google_slides_url": slides_url,
                "google_slides_id": presentation_id,
            }

        except Exception as e:
            logger.error("google_slides_update_failed", error=str(e))
            return {"status": "error", "message": str(e)}

    # ── PowerPoint builder ────────────────────────────────────────────────

    def _build_pptx(
        self,
        structure: Dict[str, Any],
        workspace_id: str,
        palette: Dict[str, RGBColor],
    ) -> Path:
        """Build a .pptx file from the generated structure."""
        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT

        slides_data = structure.get("slides", [])
        if not slides_data:
            slides_data = [
                {
                    "type": "title",
                    "title": structure.get("title", "Untitled Presentation"),
                    "subtitle": structure.get("subtitle", ""),
                }
            ]

        for i, slide_data in enumerate(slides_data):
            slide_type = slide_data.get("type", "content")
            builder = _PPTX_BUILDERS.get(slide_type, _pptx_build_content)
            builder(prs, slide_data, palette, i)

        slug = re.sub(r"[^a-z0-9]+", "-", structure.get("title", "deck").lower())[:50]
        filename = f"{workspace_id}_{slug}_{uuid4().hex[:6]}.pptx"
        file_path = self._decks_dir / filename
        prs.save(str(file_path))

        logger.info(
            "pptx_deck_created",
            path=str(file_path),
            slides=len(slides_data),
            workspace=workspace_id,
        )

        return file_path


# ══════════════════════════════════════════════════════════════════════════
# Google Slides API slide builders
# ══════════════════════════════════════════════════════════════════════════
# Each returns a list of batchUpdate request dicts for a single slide.
# Coordinates are in EMU (English Metric Units): 1 inch = 914400 EMU.

_EMU_INCH = 914400


def _gs_rgb(color: Dict[str, float]) -> Dict[str, Any]:
    """Format a color dict for the Slides API rgbColor field."""
    return {"rgbColor": color}


def _gs_text_style(
    font_size_pt: int,
    bold: bool = False,
    italic: bool = False,
    color: Optional[Dict[str, float]] = None,
) -> Dict[str, Any]:
    """Build a textStyle object for the Slides API."""
    style: Dict[str, Any] = {
        "fontSize": {"magnitude": font_size_pt, "unit": "PT"},
        "bold": bold,
        "italic": italic,
    }
    if color:
        style["foregroundColor"] = {"opaqueColor": _gs_rgb(color)}
    return style


def _gs_insert_text(object_id: str, text: str) -> Dict[str, Any]:
    return {"insertText": {"objectId": object_id, "text": text, "insertionIndex": 0}}


def _gs_create_textbox(
    slide_id: str,
    box_id: str,
    left_emu: int,
    top_emu: int,
    width_emu: int,
    height_emu: int,
) -> Dict[str, Any]:
    return {
        "createShape": {
            "objectId": box_id,
            "shapeType": "TEXT_BOX",
            "elementProperties": {
                "pageObjectId": slide_id,
                "size": {
                    "width": {"magnitude": width_emu, "unit": "EMU"},
                    "height": {"magnitude": height_emu, "unit": "EMU"},
                },
                "transform": {
                    "scaleX": 1,
                    "scaleY": 1,
                    "translateX": left_emu,
                    "translateY": top_emu,
                    "unit": "EMU",
                },
            },
        }
    }


def _gs_style_text(
    object_id: str,
    style: Dict[str, Any],
    start: int = 0,
    end: Optional[int] = None,
) -> Dict[str, Any]:
    """Apply text style to a range. If end is None, uses a special 'ALL' range."""
    text_range: Dict[str, Any]
    if end is None:
        text_range = {"type": "ALL"}
    else:
        text_range = {"type": "FIXED_RANGE", "startIndex": start, "endIndex": end}

    fields = []
    if "fontSize" in style:
        fields.append("fontSize")
    if "bold" in style:
        fields.append("bold")
    if "italic" in style:
        fields.append("italic")
    if "foregroundColor" in style:
        fields.append("foregroundColor")

    return {
        "updateTextStyle": {
            "objectId": object_id,
            "textRange": text_range,
            "style": style,
            "fields": ",".join(fields),
        }
    }


def _gs_set_bg(slide_id: str, color: Dict[str, float]) -> Dict[str, Any]:
    return {
        "updatePageProperties": {
            "objectId": slide_id,
            "pageProperties": {"pageBackgroundFill": {"solidFill": {"color": _gs_rgb(color)}}},
            "fields": "pageBackgroundFill.solidFill.color",
        }
    }


def _gs_add_speaker_notes(slide_id: str, notes: str) -> List[Dict[str, Any]]:
    """Speaker notes aren't directly settable via batchUpdate in Slides API v1.
    This is a known limitation. We skip them for Google Slides output."""
    return []


def _gs_align_text(object_id: str, alignment: str = "CENTER") -> Dict[str, Any]:
    """Return an updateParagraphStyle request to set paragraph alignment.

    alignment: "START", "CENTER", "END", "JUSTIFIED"
    """
    return {
        "updateParagraphStyle": {
            "objectId": object_id,
            "style": {"alignment": alignment},
            "textRange": {"type": "ALL"},
            "fields": "alignment",
        }
    }


def _gs_set_vertical_alignment(object_id: str, alignment: str = "MIDDLE") -> Dict[str, Any]:
    """Return an updateShapeProperties request for vertical text alignment.

    alignment: "TOP", "MIDDLE", "BOTTOM"
    """
    return {
        "updateShapeProperties": {
            "objectId": object_id,
            "shapeProperties": {
                "contentAlignment": alignment,
            },
            "fields": "contentAlignment",
        }
    }


def _gs_build_title(
    slide_id: str, data: Dict[str, Any], pal: Dict[str, Dict[str, float]], idx: int
) -> List[Dict[str, Any]]:
    reqs: List[Dict[str, Any]] = []

    # Title textbox — centered on 10" x 5.625" canvas
    title_id = f"{slide_id}_title"
    title_text = data.get("title", "")
    reqs.append(
        _gs_create_textbox(
            slide_id,
            title_id,
            left_emu=int(0.6 * _EMU_INCH),
            top_emu=int(1.2 * _EMU_INCH),
            width_emu=int(8.8 * _EMU_INCH),
            height_emu=int(1.5 * _EMU_INCH),
        )
    )
    if title_text:
        reqs.append(_gs_insert_text(title_id, title_text))
        reqs.append(_gs_style_text(title_id, _gs_text_style(44, bold=True, color=pal["dark"])))
    reqs.append(_gs_align_text(title_id, "CENTER"))
    reqs.append(_gs_set_vertical_alignment(title_id, "MIDDLE"))

    # Subtitle
    subtitle_text = data.get("subtitle", "")
    if subtitle_text:
        sub_id = f"{slide_id}_subtitle"
        reqs.append(
            _gs_create_textbox(
                slide_id,
                sub_id,
                left_emu=int(0.6 * _EMU_INCH),
                top_emu=int(2.9 * _EMU_INCH),
                width_emu=int(8.8 * _EMU_INCH),
                height_emu=int(1.0 * _EMU_INCH),
            )
        )
        reqs.append(_gs_insert_text(sub_id, subtitle_text))
        reqs.append(_gs_style_text(sub_id, _gs_text_style(22, color=pal["muted"])))
        reqs.append(_gs_align_text(sub_id, "CENTER"))

    return reqs


def _gs_build_section(
    slide_id: str, data: Dict[str, Any], pal: Dict[str, Dict[str, float]], idx: int
) -> List[Dict[str, Any]]:
    reqs: List[Dict[str, Any]] = []

    title_id = f"{slide_id}_title"
    title_text = data.get("title", "")
    reqs.append(
        _gs_create_textbox(
            slide_id,
            title_id,
            left_emu=int(0.6 * _EMU_INCH),
            top_emu=int(1.5 * _EMU_INCH),
            width_emu=int(8.8 * _EMU_INCH),
            height_emu=int(2.0 * _EMU_INCH),
        )
    )
    text = title_text
    subtitle = data.get("subtitle", "")
    if subtitle:
        text += f"\n{subtitle}"
    if text:
        reqs.append(_gs_insert_text(title_id, text))
        reqs.append(_gs_style_text(title_id, _gs_text_style(40, bold=True, color=pal["primary"])))
    reqs.append(_gs_align_text(title_id, "CENTER"))
    reqs.append(_gs_set_vertical_alignment(title_id, "MIDDLE"))

    return reqs


def _gs_build_content(
    slide_id: str, data: Dict[str, Any], pal: Dict[str, Dict[str, float]], idx: int
) -> List[Dict[str, Any]]:
    reqs: List[Dict[str, Any]] = []

    # Title
    title_id = f"{slide_id}_title"
    title_text = data.get("title", "")
    reqs.append(
        _gs_create_textbox(
            slide_id,
            title_id,
            left_emu=int(0.6 * _EMU_INCH),
            top_emu=int(0.4 * _EMU_INCH),
            width_emu=int(8.8 * _EMU_INCH),
            height_emu=int(0.7 * _EMU_INCH),
        )
    )
    if title_text:
        reqs.append(_gs_insert_text(title_id, title_text))
        reqs.append(_gs_style_text(title_id, _gs_text_style(32, bold=True, color=pal["dark"])))
    reqs.append(_gs_align_text(title_id, "START"))

    # Bullets
    bullets = data.get("bullets", data.get("content", []))
    if isinstance(bullets, str):
        bullets = [b.strip() for b in bullets.split("\n") if b.strip()]

    if bullets:
        body_id = f"{slide_id}_body"
        body_text = "\n".join(f"• {b}" for b in bullets)
        reqs.append(
            _gs_create_textbox(
                slide_id,
                body_id,
                left_emu=int(0.6 * _EMU_INCH),
                top_emu=int(1.3 * _EMU_INCH),
                width_emu=int(8.4 * _EMU_INCH),
                height_emu=int(3.8 * _EMU_INCH),
            )
        )
        reqs.append(_gs_insert_text(body_id, body_text))
        reqs.append(_gs_style_text(body_id, _gs_text_style(20, color=pal["text"])))
        reqs.append(_gs_align_text(body_id, "START"))
        reqs.append(_gs_set_vertical_alignment(body_id, "TOP"))

    return reqs


def _gs_build_two_column(
    slide_id: str, data: Dict[str, Any], pal: Dict[str, Dict[str, float]], idx: int
) -> List[Dict[str, Any]]:
    reqs: List[Dict[str, Any]] = []

    # Title
    title_id = f"{slide_id}_title"
    title_text = data.get("title", "")
    reqs.append(
        _gs_create_textbox(
            slide_id,
            title_id,
            left_emu=int(0.6 * _EMU_INCH),
            top_emu=int(0.4 * _EMU_INCH),
            width_emu=int(8.8 * _EMU_INCH),
            height_emu=int(0.7 * _EMU_INCH),
        )
    )
    if title_text:
        reqs.append(_gs_insert_text(title_id, title_text))
        reqs.append(_gs_style_text(title_id, _gs_text_style(32, bold=True, color=pal["dark"])))
    reqs.append(_gs_align_text(title_id, "START"))

    # Left column
    left_title = data.get("left_title", data.get("column_1_title", ""))
    left_items = data.get("left_column", data.get("column_1", []))
    if isinstance(left_items, list):
        left_items = "\n".join(f"• {i}" for i in left_items)
    left_text = f"{left_title}\n\n{left_items}" if left_title else str(left_items)
    left_id = f"{slide_id}_left"
    reqs.append(
        _gs_create_textbox(
            slide_id,
            left_id,
            left_emu=int(0.6 * _EMU_INCH),
            top_emu=int(1.3 * _EMU_INCH),
            width_emu=int(4.0 * _EMU_INCH),
            height_emu=int(3.5 * _EMU_INCH),
        )
    )
    if left_text.strip():
        reqs.append(_gs_insert_text(left_id, left_text))
        reqs.append(_gs_style_text(left_id, _gs_text_style(18, color=pal["text"])))
    reqs.append(_gs_align_text(left_id, "START"))
    reqs.append(_gs_set_vertical_alignment(left_id, "TOP"))

    # Right column
    right_title = data.get("right_title", data.get("column_2_title", ""))
    right_items = data.get("right_column", data.get("column_2", []))
    if isinstance(right_items, list):
        right_items = "\n".join(f"• {i}" for i in right_items)
    right_text = f"{right_title}\n\n{right_items}" if right_title else str(right_items)
    right_id = f"{slide_id}_right"
    reqs.append(
        _gs_create_textbox(
            slide_id,
            right_id,
            left_emu=int(5.2 * _EMU_INCH),
            top_emu=int(1.3 * _EMU_INCH),
            width_emu=int(4.0 * _EMU_INCH),
            height_emu=int(3.5 * _EMU_INCH),
        )
    )
    if right_text.strip():
        reqs.append(_gs_insert_text(right_id, right_text))
        reqs.append(_gs_style_text(right_id, _gs_text_style(18, color=pal["text"])))
    reqs.append(_gs_align_text(right_id, "START"))
    reqs.append(_gs_set_vertical_alignment(right_id, "TOP"))

    return reqs


def _gs_build_stats(
    slide_id: str, data: Dict[str, Any], pal: Dict[str, Dict[str, float]], idx: int
) -> List[Dict[str, Any]]:
    reqs: List[Dict[str, Any]] = []

    # Title
    title_id = f"{slide_id}_title"
    title_text = data.get("title", "Key Metrics")
    reqs.append(
        _gs_create_textbox(
            slide_id,
            title_id,
            left_emu=int(0.6 * _EMU_INCH),
            top_emu=int(0.4 * _EMU_INCH),
            width_emu=int(8.8 * _EMU_INCH),
            height_emu=int(0.7 * _EMU_INCH),
        )
    )
    if title_text:
        reqs.append(_gs_insert_text(title_id, title_text))
        reqs.append(_gs_style_text(title_id, _gs_text_style(32, bold=True, color=pal["dark"])))
    reqs.append(_gs_align_text(title_id, "START"))

    # Stats as big numbers — fit within 10" page
    stats = data.get("stats", data.get("metrics", []))
    if isinstance(stats, list) and stats:
        num_stats = len(stats)
        # Available width: 10" - 2*0.6" margin = 8.8", each stat box 2.2" wide
        stat_w = 2.2
        total_stats_width = num_stats * stat_w
        # Center stats horizontally if they fit, otherwise compress spacing
        if total_stats_width <= 8.8:
            start_left = (_GS_W - total_stats_width) / 2
            spacing = stat_w
        else:
            start_left = _GS_MX
            spacing = min(stat_w, (8.8) / num_stats)
        for j, stat in enumerate(stats):
            val = str(stat.get("value", stat.get("number", "")))
            label = str(stat.get("label", stat.get("description", "")))
            text = f"{val}\n{label}"
            stat_id = f"{slide_id}_stat{j}"
            reqs.append(
                _gs_create_textbox(
                    slide_id,
                    stat_id,
                    left_emu=int((start_left + j * spacing) * _EMU_INCH),
                    top_emu=int(1.8 * _EMU_INCH),
                    width_emu=int(stat_w * _EMU_INCH),
                    height_emu=int(2.0 * _EMU_INCH),
                )
            )
            reqs.append(_gs_insert_text(stat_id, text))
            # Style the value part large, label part smaller
            val_len = len(val)
            reqs.append(
                _gs_style_text(
                    stat_id,
                    _gs_text_style(48, bold=True, color=pal["accent"]),
                    start=0,
                    end=val_len,
                )
            )
            if label:
                reqs.append(
                    _gs_style_text(
                        stat_id,
                        _gs_text_style(16, color=pal["muted"]),
                        start=val_len + 1,
                        end=val_len + 1 + len(label),
                    )
                )
            reqs.append(_gs_align_text(stat_id, "CENTER"))
            reqs.append(_gs_set_vertical_alignment(stat_id, "MIDDLE"))

    return reqs


def _gs_build_quote(
    slide_id: str, data: Dict[str, Any], pal: Dict[str, Dict[str, float]], idx: int
) -> List[Dict[str, Any]]:
    reqs: List[Dict[str, Any]] = []

    quote = data.get("quote", data.get("text", ""))
    attribution = data.get("attribution", data.get("author", ""))

    quote_id = f"{slide_id}_quote"
    full_text = f"\u201c{quote}\u201d"
    if attribution:
        full_text += f"\n\n\u2014 {attribution}"

    reqs.append(
        _gs_create_textbox(
            slide_id,
            quote_id,
            left_emu=int(1.0 * _EMU_INCH),
            top_emu=int(1.0 * _EMU_INCH),
            width_emu=int(8.0 * _EMU_INCH),
            height_emu=int(3.0 * _EMU_INCH),
        )
    )
    if full_text:
        reqs.append(_gs_insert_text(quote_id, full_text))
        reqs.append(_gs_style_text(quote_id, _gs_text_style(28, italic=True, color=pal["text"])))
    reqs.append(_gs_align_text(quote_id, "CENTER"))
    reqs.append(_gs_set_vertical_alignment(quote_id, "MIDDLE"))

    return reqs


def _gs_build_image(
    slide_id: str, data: Dict[str, Any], pal: Dict[str, Dict[str, float]], idx: int
) -> List[Dict[str, Any]]:
    """Image slide — inserts image from URL if provided, otherwise placeholder text."""
    reqs: List[Dict[str, Any]] = []

    title_text = data.get("title", "")
    caption = data.get("caption", data.get("description", ""))
    image_url = data.get("image_url", "")

    if title_text:
        title_id = f"{slide_id}_title"
        reqs.append(
            _gs_create_textbox(
                slide_id,
                title_id,
                left_emu=int(0.6 * _EMU_INCH),
                top_emu=int(0.4 * _EMU_INCH),
                width_emu=int(8.8 * _EMU_INCH),
                height_emu=int(0.7 * _EMU_INCH),
            )
        )
        reqs.append(_gs_insert_text(title_id, title_text))
        reqs.append(_gs_style_text(title_id, _gs_text_style(32, bold=True, color=pal["dark"])))
        reqs.append(_gs_align_text(title_id, "START"))

    if image_url:
        img_id = f"{slide_id}_img"
        reqs.append(
            {
                "createImage": {
                    "objectId": img_id,
                    "url": image_url,
                    "elementProperties": {
                        "pageObjectId": slide_id,
                        "size": {
                            "width": {"magnitude": int(8.0 * _EMU_INCH), "unit": "EMU"},
                            "height": {"magnitude": int(3.2 * _EMU_INCH), "unit": "EMU"},
                        },
                        "transform": {
                            "scaleX": 1,
                            "scaleY": 1,
                            "translateX": int(1.0 * _EMU_INCH),
                            "translateY": int(1.3 * _EMU_INCH),
                            "unit": "EMU",
                        },
                    },
                }
            }
        )
    else:
        placeholder_id = f"{slide_id}_imgplaceholder"
        reqs.append(
            _gs_create_textbox(
                slide_id,
                placeholder_id,
                left_emu=int(1.0 * _EMU_INCH),
                top_emu=int(1.3 * _EMU_INCH),
                width_emu=int(8.0 * _EMU_INCH),
                height_emu=int(3.2 * _EMU_INCH),
            )
        )
        reqs.append(_gs_insert_text(placeholder_id, "[Image Placeholder]"))
        reqs.append(_gs_style_text(placeholder_id, _gs_text_style(24, color=pal["muted"])))
        reqs.append(_gs_align_text(placeholder_id, "CENTER"))
        reqs.append(_gs_set_vertical_alignment(placeholder_id, "MIDDLE"))

    if caption:
        cap_id = f"{slide_id}_caption"
        reqs.append(
            _gs_create_textbox(
                slide_id,
                cap_id,
                left_emu=int(1.0 * _EMU_INCH),
                top_emu=int(4.7 * _EMU_INCH),
                width_emu=int(8.0 * _EMU_INCH),
                height_emu=int(0.6 * _EMU_INCH),
            )
        )
        reqs.append(_gs_insert_text(cap_id, caption))
        reqs.append(_gs_style_text(cap_id, _gs_text_style(16, italic=True, color=pal["muted"])))
        reqs.append(_gs_align_text(cap_id, "CENTER"))

    return reqs


def _gs_build_cta(
    slide_id: str, data: Dict[str, Any], pal: Dict[str, Dict[str, float]], idx: int
) -> List[Dict[str, Any]]:
    reqs: List[Dict[str, Any]] = []

    cta_text = data.get("title", data.get("cta", "Let's Get Started"))
    cta_id = f"{slide_id}_cta"
    reqs.append(
        _gs_create_textbox(
            slide_id,
            cta_id,
            left_emu=int(0.6 * _EMU_INCH),
            top_emu=int(1.2 * _EMU_INCH),
            width_emu=int(8.8 * _EMU_INCH),
            height_emu=int(1.5 * _EMU_INCH),
        )
    )
    if cta_text:
        reqs.append(_gs_insert_text(cta_id, cta_text))
        reqs.append(_gs_style_text(cta_id, _gs_text_style(44, bold=True, color=pal["primary"])))
    reqs.append(_gs_align_text(cta_id, "CENTER"))
    reqs.append(_gs_set_vertical_alignment(cta_id, "MIDDLE"))

    details = data.get("subtitle", data.get("details", data.get("contact", "")))
    if details:
        det_id = f"{slide_id}_details"
        reqs.append(
            _gs_create_textbox(
                slide_id,
                det_id,
                left_emu=int(1.2 * _EMU_INCH),
                top_emu=int(3.2 * _EMU_INCH),
                width_emu=int(7.0 * _EMU_INCH),
                height_emu=int(1.5 * _EMU_INCH),
            )
        )
        reqs.append(_gs_insert_text(det_id, details))
        reqs.append(_gs_style_text(det_id, _gs_text_style(22, color=pal["dark"])))
        reqs.append(_gs_align_text(det_id, "CENTER"))

    return reqs


_GS_SLIDE_BUILDERS = {
    "title": _gs_build_title,
    "section": _gs_build_section,
    "content": _gs_build_content,
    "two_column": _gs_build_two_column,
    "stats": _gs_build_stats,
    "quote": _gs_build_quote,
    "image": _gs_build_image,
    "cta": _gs_build_cta,
}


# ══════════════════════════════════════════════════════════════════════════
# PowerPoint (.pptx) slide builders
# ══════════════════════════════════════════════════════════════════════════


def _add_blank_slide(prs: Presentation) -> Any:
    """Add a blank slide and return its slide object."""
    layout = prs.slide_layouts[6]  # Blank layout
    return prs.slides.add_slide(layout)


def _pptx_build_title(
    prs: Presentation, data: Dict[str, Any], palette: Dict[str, RGBColor], idx: int
) -> None:
    slide = _add_blank_slide(prs)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = palette["dark"]

    title_text = data.get("title", "")
    left = Inches(1)
    txBox = slide.shapes.add_textbox(left, Inches(2.0), Inches(11.3), Inches(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = palette["text_light"]
    p.alignment = PP_ALIGN.LEFT

    subtitle_text = data.get("subtitle", "")
    if subtitle_text:
        txBox2 = slide.shapes.add_textbox(left, Inches(4.2), Inches(11.3), Inches(1.5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle_text
        p2.font.size = Pt(22)
        p2.font.color.rgb = palette["muted"]
        p2.alignment = PP_ALIGN.LEFT

    shape = slide.shapes.add_shape(1, left, Inches(4.0), Inches(2), Inches(0.06))
    shape.fill.solid()
    shape.fill.fore_color.rgb = palette["primary"]
    shape.line.fill.background()

    _pptx_add_speaker_notes(slide, data)


def _pptx_build_section(
    prs: Presentation, data: Dict[str, Any], palette: Dict[str, RGBColor], idx: int
) -> None:
    slide = _add_blank_slide(prs)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = palette["primary"]

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data.get("title", "")
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = palette["text_light"]
    p.alignment = PP_ALIGN.LEFT

    subtitle = data.get("subtitle", "")
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p2.space_before = Pt(12)

    _pptx_add_speaker_notes(slide, data)


def _pptx_build_content(
    prs: Presentation, data: Dict[str, Any], palette: Dict[str, RGBColor], idx: int
) -> None:
    slide = _add_blank_slide(prs)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = palette["light"]

    left = Inches(0.8)
    txBox = slide.shapes.add_textbox(left, Inches(0.5), Inches(11.7), Inches(1.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data.get("title", "")
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = palette["dark"]

    bar = slide.shapes.add_shape(1, left, Inches(1.4), Inches(1.5), Inches(0.04))
    bar.fill.solid()
    bar.fill.fore_color.rgb = palette["primary"]
    bar.line.fill.background()

    bullets = data.get("bullets", data.get("content", []))
    if isinstance(bullets, str):
        bullets = [b.strip() for b in bullets.split("\n") if b.strip()]

    if bullets:
        txBox2 = slide.shapes.add_textbox(left, Inches(1.8), Inches(11.0), Inches(5.0))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        for j, bullet in enumerate(bullets):
            p_b = tf2.paragraphs[0] if j == 0 else tf2.add_paragraph()
            p_b.text = bullet
            p_b.font.size = Pt(20)
            p_b.font.color.rgb = palette["text"]
            p_b.space_before = Pt(10)
            p_b.level = 0

    _pptx_add_speaker_notes(slide, data)


def _pptx_build_two_column(
    prs: Presentation, data: Dict[str, Any], palette: Dict[str, RGBColor], idx: int
) -> None:
    slide = _add_blank_slide(prs)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = palette["light"]

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data.get("title", "")
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = palette["dark"]

    _pptx_add_column(
        slide,
        Inches(0.8),
        data.get("left_title", data.get("column_1_title", "")),
        data.get("left_column", data.get("column_1", [])),
        palette,
    )
    _pptx_add_column(
        slide,
        Inches(7.0),
        data.get("right_title", data.get("column_2_title", "")),
        data.get("right_column", data.get("column_2", [])),
        palette,
    )

    _pptx_add_speaker_notes(slide, data)


def _pptx_build_stats(
    prs: Presentation, data: Dict[str, Any], palette: Dict[str, RGBColor], idx: int
) -> None:
    slide = _add_blank_slide(prs)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = palette["dark"]

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data.get("title", "Key Metrics")
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = palette["text_light"]

    stats = data.get("stats", data.get("metrics", []))
    if isinstance(stats, list) and stats:
        num_stats = len(stats)
        spacing = min(3.5, 10.0 / num_stats)
        for j, stat in enumerate(stats):
            col_left = Inches(0.8 + j * spacing)
            val = str(stat.get("value", stat.get("number", "")))
            txVal = slide.shapes.add_textbox(col_left, Inches(2.5), Inches(3.0), Inches(1.5))
            pVal = txVal.text_frame.paragraphs[0]
            pVal.text = val
            pVal.font.size = Pt(48)
            pVal.font.bold = True
            pVal.font.color.rgb = palette["accent"]
            pVal.alignment = PP_ALIGN.CENTER

            label = str(stat.get("label", stat.get("description", "")))
            txLbl = slide.shapes.add_textbox(col_left, Inches(4.0), Inches(3.0), Inches(1.0))
            txLbl.text_frame.word_wrap = True
            pLbl = txLbl.text_frame.paragraphs[0]
            pLbl.text = label
            pLbl.font.size = Pt(16)
            pLbl.font.color.rgb = palette["muted"]
            pLbl.alignment = PP_ALIGN.CENTER

    _pptx_add_speaker_notes(slide, data)


def _pptx_build_quote(
    prs: Presentation, data: Dict[str, Any], palette: Dict[str, RGBColor], idx: int
) -> None:
    slide = _add_blank_slide(prs)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = palette["light"]

    txQ = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(2), Inches(2))
    pQ = txQ.text_frame.paragraphs[0]
    pQ.text = "\u201c"
    pQ.font.size = Pt(120)
    pQ.font.color.rgb = palette["primary"]

    quote = data.get("quote", data.get("text", ""))
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.0), Inches(3.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = quote
    p.font.size = Pt(28)
    p.font.italic = True
    p.font.color.rgb = palette["text"]

    attribution = data.get("attribution", data.get("author", ""))
    if attribution:
        txA = slide.shapes.add_textbox(Inches(1.5), Inches(5.5), Inches(10.0), Inches(0.8))
        pA = txA.text_frame.paragraphs[0]
        pA.text = f"\u2014 {attribution}"
        pA.font.size = Pt(18)
        pA.font.color.rgb = palette["muted"]

    _pptx_add_speaker_notes(slide, data)


def _pptx_build_image(
    prs: Presentation, data: Dict[str, Any], palette: Dict[str, RGBColor], idx: int
) -> None:
    slide = _add_blank_slide(prs)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = palette["light"]

    title_text = data.get("title", "")
    caption = data.get("caption", data.get("description", ""))
    image_url = data.get("image_url", "")

    if title_text:
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = palette["dark"]

    img_left = Inches(1.5)
    img_top = Inches(1.8)
    img_width = Inches(10.3)
    img_height = Inches(4.5)

    if image_url:
        try:
            import urllib.request

            with urllib.request.urlopen(image_url, timeout=15) as resp:
                img_data = resp.read()
            slide.shapes.add_picture(BytesIO(img_data), img_left, img_top, img_width, img_height)
        except Exception as e:
            logger.warning("image_download_failed", url=image_url, error=str(e))
            _pptx_add_image_placeholder(slide, img_left, img_top, img_width, img_height, palette)
    else:
        _pptx_add_image_placeholder(slide, img_left, img_top, img_width, img_height, palette)

    if caption:
        txC = slide.shapes.add_textbox(Inches(1.5), Inches(6.5), Inches(10.3), Inches(0.8))
        txC.text_frame.word_wrap = True
        pC = txC.text_frame.paragraphs[0]
        pC.text = caption
        pC.font.size = Pt(16)
        pC.font.italic = True
        pC.font.color.rgb = palette["muted"]
        pC.alignment = PP_ALIGN.CENTER

    _pptx_add_speaker_notes(slide, data)


def _pptx_build_cta(
    prs: Presentation, data: Dict[str, Any], palette: Dict[str, RGBColor], idx: int
) -> None:
    slide = _add_blank_slide(prs)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = palette["primary"]

    cta_text = data.get("title", data.get("cta", "Let's Get Started"))
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11.3), Inches(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = cta_text
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = palette["text_light"]
    p.alignment = PP_ALIGN.CENTER

    details = data.get("subtitle", data.get("details", data.get("contact", "")))
    if details:
        txD = slide.shapes.add_textbox(Inches(2), Inches(4.5), Inches(9.3), Inches(2.0))
        txD.text_frame.word_wrap = True
        pD = txD.text_frame.paragraphs[0]
        pD.text = details
        pD.font.size = Pt(22)
        pD.font.color.rgb = palette["text_light"]
        pD.alignment = PP_ALIGN.CENTER

    _pptx_add_speaker_notes(slide, data)


_PPTX_BUILDERS = {
    "title": _pptx_build_title,
    "section": _pptx_build_section,
    "content": _pptx_build_content,
    "two_column": _pptx_build_two_column,
    "stats": _pptx_build_stats,
    "quote": _pptx_build_quote,
    "image": _pptx_build_image,
    "cta": _pptx_build_cta,
}


# ══════════════════════════════════════════════════════════════════════════
# Helpers
# ══════════════════════════════════════════════════════════════════════════


def _get_palette(workspace_id: str) -> Dict[str, RGBColor]:
    return BRAND_PALETTES.get(workspace_id, DEFAULT_PALETTE)


def _palette_to_floats(palette: Dict[str, RGBColor]) -> Dict[str, Dict[str, float]]:
    """Convert RGBColor palette to float dicts (0.0-1.0) for Google Slides API."""
    result = {}
    for key, color in palette.items():
        try:
            result[key] = {
                "red": color.red / 255.0,
                "green": color.green / 255.0,
                "blue": color.blue / 255.0,
            }
        except (AttributeError, TypeError):
            result[key] = {"red": 0.2, "green": 0.2, "blue": 0.2}
    return result


def _pptx_add_speaker_notes(slide: Any, data: Dict[str, Any]) -> None:
    notes = data.get("speaker_notes", "")
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def _pptx_add_column(
    slide: Any, left: Any, title: str, items: Any, palette: Dict[str, RGBColor]
) -> None:
    col_width = Inches(5.5)
    if title:
        txT = slide.shapes.add_textbox(left, Inches(1.8), col_width, Inches(0.6))
        pT = txT.text_frame.paragraphs[0]
        pT.text = title
        pT.font.size = Pt(22)
        pT.font.bold = True
        pT.font.color.rgb = palette["primary"]

    if isinstance(items, str):
        items = [i.strip() for i in items.split("\n") if i.strip()]
    if isinstance(items, list) and items:
        txI = slide.shapes.add_textbox(left, Inches(2.5), col_width, Inches(4.0))
        tfI = txI.text_frame
        tfI.word_wrap = True
        for j, item in enumerate(items):
            p_i = tfI.paragraphs[0] if j == 0 else tfI.add_paragraph()
            p_i.text = str(item)
            p_i.font.size = Pt(18)
            p_i.font.color.rgb = palette["text"]
            p_i.space_before = Pt(8)


def _pptx_add_image_placeholder(
    slide: Any, left: Any, top: Any, width: Any, height: Any, palette: Dict[str, RGBColor]
) -> None:
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
    shape.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = "[Image Placeholder]"
    p.font.size = Pt(18)
    p.font.color.rgb = palette["muted"]
    tf.paragraphs[0].space_before = Pt(80)


def _structure_format_instructions(extra_fields: str = "") -> str:
    return f"""Return ONLY a valid JSON object with this structure:
{{
  "title": "Deck Title",
  "subtitle": "Optional subtitle",
  "slides": [
    {{
      "type": "title",
      "title": "Slide Title",
      "subtitle": "Slide subtitle"
      {", " + extra_fields if extra_fields else ""}
    }},
    {{
      "type": "content",
      "title": "Slide Title",
      "bullets": ["Point 1", "Point 2", "Point 3"]
      {", " + extra_fields if extra_fields else ""}
    }},
    {{
      "type": "two_column",
      "title": "Comparison",
      "left_title": "Option A",
      "left_column": ["Point 1", "Point 2"],
      "right_title": "Option B",
      "right_column": ["Point 1", "Point 2"]
    }},
    {{
      "type": "stats",
      "title": "Key Metrics",
      "stats": [{{"value": "85%", "label": "Growth rate"}}]
    }},
    {{
      "type": "quote",
      "quote": "Quote text",
      "attribution": "Author Name"
    }},
    {{
      "type": "image",
      "title": "Visual Title",
      "caption": "Image description"
    }},
    {{
      "type": "section",
      "title": "Section Name"
    }},
    {{
      "type": "cta",
      "title": "Call to Action",
      "subtitle": "Contact info or next steps"
    }}
  ]
}}

Use a variety of slide types for visual interest. Ensure 'type' is one of:
title, section, content, two_column, stats, quote, image, cta"""


def _parse_structure(raw: str) -> Dict[str, Any]:
    try:
        return json.loads(raw)
    except (json.JSONDecodeError, TypeError):
        pass

    match = re.search(r"```(?:json)?\s*\n?(.*?)\n?```", raw, re.DOTALL)
    if match:
        try:
            return json.loads(match.group(1))
        except (json.JSONDecodeError, TypeError):
            pass

    logger.warning("deck_structure_parse_fallback", raw_length=len(raw))
    lines = [l.strip() for l in raw.split("\n") if l.strip()]
    return {
        "title": lines[0] if lines else "Untitled Presentation",
        "slides": [
            {"type": "title", "title": lines[0] if lines else "Untitled"},
            {
                "type": "content",
                "title": "Overview",
                "bullets": lines[1:8] if len(lines) > 1 else ["Content to be added"],
            },
            {"type": "cta", "title": "Thank You", "subtitle": "Questions?"},
        ],
    }
