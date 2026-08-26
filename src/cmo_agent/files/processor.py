"""File processor — downloads Slack files and extracts content for LLM consumption.

Supports:
- Images (PNG, JPG, GIF, WEBP) → sent as vision content blocks to Claude
- PDFs → text extraction via pypdf
- Spreadsheets (XLSX, CSV) → cell data extraction via openpyxl/csv
- Word documents (DOCX) → text extraction via python-docx
- PowerPoint (PPTX) → text extraction via python-pptx
- Plain text / code / markdown → raw text
"""

from __future__ import annotations

import base64
import csv
import io
import mimetypes
from typing import Any, Dict, List, Optional, Tuple

import httpx
import structlog
from PIL import Image

logger = structlog.get_logger()

# File type categories
IMAGE_MIMES = {
    "image/png",
    "image/jpeg",
    "image/gif",
    "image/webp",
}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".webp"}

PDF_MIMES = {"application/pdf"}

SPREADSHEET_MIMES = {
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/vnd.ms-excel",
    "text/csv",
}
SPREADSHEET_EXTENSIONS = {".xlsx", ".xls", ".csv"}

DOCX_MIMES = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/msword",
}
DOCX_EXTENSIONS = {".docx", ".doc"}

PPTX_MIMES = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/vnd.ms-powerpoint",
}
PPTX_EXTENSIONS = {".pptx", ".ppt"}

TEXT_MIMES = {
    "text/plain",
    "text/markdown",
    "text/html",
    "text/css",
    "application/json",
    "application/javascript",
    "application/xml",
    "text/x-python",
    "text/x-java",
    "text/x-c",
}
TEXT_EXTENSIONS = {
    ".txt",
    ".md",
    ".html",
    ".css",
    ".json",
    ".js",
    ".ts",
    ".py",
    ".java",
    ".c",
    ".cpp",
    ".rb",
    ".go",
    ".rs",
    ".yaml",
    ".yml",
    ".xml",
    ".toml",
    ".ini",
    ".cfg",
    ".sh",
    ".bash",
    ".sql",
}

# Maximum text extraction size (chars) to avoid context overflow
MAX_TEXT_CHARS = 50_000
MAX_SPREADSHEET_ROWS = 500

# Anthropic base64 image limit is 5MB; target below that with headroom
MAX_IMAGE_B64_BYTES = 4_500_000  # ~4.5MB to leave margin under 5MB API limit


class FileAttachment:
    """Represents a processed file ready for LLM consumption."""

    def __init__(
        self,
        filename: str,
        file_type: str,  # "image", "text", "pdf", "spreadsheet", "document", "presentation", "unknown"
        mime_type: str,
        content_blocks: List[Dict[str, Any]],  # Anthropic content blocks
        summary: str = "",  # Brief description for logging
        size_bytes: int = 0,
    ):
        self.filename = filename
        self.file_type = file_type
        self.mime_type = mime_type
        self.content_blocks = content_blocks
        self.summary = summary
        self.size_bytes = size_bytes


def _compress_image_for_api(data: bytes, media_type: str) -> Tuple[bytes, str]:
    """Compress/resize an image so its base64 encoding stays under the API limit.

    Returns (image_bytes, media_type) — may convert non-JPEG to JPEG for
    better compression.  GIFs are kept as-is (only first frame would survive
    a PIL re-encode, which isn't useful).
    """
    b64_len = len(data) * 4 // 3  # base64 is ~4/3 of raw bytes
    if b64_len <= MAX_IMAGE_B64_BYTES:
        return data, media_type

    # GIFs can't be meaningfully resized without losing animation
    if media_type == "image/gif":
        logger.warning("gif_too_large", raw_bytes=len(data))
        return data, media_type

    img = Image.open(io.BytesIO(data))
    if img.mode in ("RGBA", "P"):
        img = img.convert("RGB")
    out_mime = "image/jpeg"

    # Iteratively downscale until the base64 size fits
    quality = 85
    for scale in (1.0, 0.75, 0.5, 0.35, 0.25):
        w = int(img.width * scale)
        h = int(img.height * scale)
        resized = img.resize((w, h), Image.LANCZOS) if scale < 1.0 else img

        buf = io.BytesIO()
        resized.save(buf, format="JPEG", quality=quality, optimize=True)
        raw = buf.getvalue()
        if len(raw) * 4 // 3 <= MAX_IMAGE_B64_BYTES:
            logger.info(
                "image_compressed",
                original_bytes=len(data),
                compressed_bytes=len(raw),
                scale=scale,
                quality=quality,
                dimensions=f"{w}x{h}",
            )
            return raw, out_mime

        # Try lower quality at same scale before going smaller
        if quality > 60:
            buf2 = io.BytesIO()
            resized.save(buf2, format="JPEG", quality=60, optimize=True)
            raw2 = buf2.getvalue()
            if len(raw2) * 4 // 3 <= MAX_IMAGE_B64_BYTES:
                logger.info(
                    "image_compressed",
                    original_bytes=len(data),
                    compressed_bytes=len(raw2),
                    scale=scale,
                    quality=60,
                    dimensions=f"{w}x{h}",
                )
                return raw2, out_mime

    # Last resort — return whatever we have at smallest scale
    logger.warning("image_compression_fallback", final_bytes=len(raw))
    return raw, out_mime


def classify_file(filename: str, mimetype: Optional[str] = None) -> str:
    """Classify a file into a processing category."""
    ext = "." + filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    mime = mimetype or mimetypes.guess_type(filename)[0] or ""

    if mime in IMAGE_MIMES or ext in IMAGE_EXTENSIONS:
        return "image"
    if mime in PDF_MIMES or ext == ".pdf":
        return "pdf"
    if mime in SPREADSHEET_MIMES or ext in SPREADSHEET_EXTENSIONS:
        return "spreadsheet"
    if mime in DOCX_MIMES or ext in DOCX_EXTENSIONS:
        return "document"
    if mime in PPTX_MIMES or ext in PPTX_EXTENSIONS:
        return "presentation"
    if mime in TEXT_MIMES or ext in TEXT_EXTENSIONS:
        return "text"
    return "unknown"


async def download_slack_file(
    file_url: str, bot_token: str, max_size: int = 20 * 1024 * 1024
) -> Tuple[bytes, str]:
    """Download a file from Slack's servers.

    Args:
        file_url: The url_private from the Slack file object
        bot_token: Slack bot token for authorization
        max_size: Maximum file size in bytes (default 20MB)

    Returns:
        (file_bytes, content_type)
    """
    async with httpx.AsyncClient(timeout=30.0) as client:
        response = await client.get(
            file_url,
            headers={"Authorization": f"Bearer {bot_token}"},
        )
        response.raise_for_status()

        content_type = response.headers.get("content-type", "application/octet-stream")
        data = response.content

        if len(data) > max_size:
            raise ValueError(f"File too large: {len(data)} bytes (max {max_size})")

        return data, content_type


def extract_text_from_pdf(data: bytes) -> str:
    """Extract text from a PDF file."""
    try:
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(data))
        pages = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            if text.strip():
                pages.append(f"--- Page {i + 1} ---\n{text}")
        return "\n\n".join(pages)[:MAX_TEXT_CHARS]
    except Exception as e:
        logger.warning("pdf_extraction_failed", error=str(e))
        return f"[PDF extraction failed: {e}]"


def extract_text_from_docx(data: bytes) -> str:
    """Extract text from a Word document."""
    try:
        from docx import Document

        doc = Document(io.BytesIO(data))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(paragraphs)[:MAX_TEXT_CHARS]
    except Exception as e:
        logger.warning("docx_extraction_failed", error=str(e))
        return f"[DOCX extraction failed: {e}]"


def extract_text_from_pptx(data: bytes) -> str:
    """Extract text from a PowerPoint presentation."""
    try:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(data))
        slides = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            texts.append(text)
            if texts:
                slides.append(f"--- Slide {i + 1} ---\n" + "\n".join(texts))
        return "\n\n".join(slides)[:MAX_TEXT_CHARS]
    except Exception as e:
        logger.warning("pptx_extraction_failed", error=str(e))
        return f"[PPTX extraction failed: {e}]"


def extract_text_from_spreadsheet(data: bytes, filename: str) -> str:
    """Extract text from a spreadsheet (XLSX or CSV)."""
    ext = "." + filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    try:
        if ext == ".csv":
            text = data.decode("utf-8", errors="replace")
            reader = csv.reader(io.StringIO(text))
            rows = []
            for i, row in enumerate(reader):
                if i >= MAX_SPREADSHEET_ROWS:
                    rows.append(f"... ({i} rows total, truncated)")
                    break
                rows.append(" | ".join(str(cell) for cell in row))
            return "\n".join(rows)[:MAX_TEXT_CHARS]
        else:
            from openpyxl import load_workbook

            wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
            sheets = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = []
                for i, row in enumerate(ws.iter_rows(values_only=True)):
                    if i >= MAX_SPREADSHEET_ROWS:
                        rows.append(f"... (truncated at {MAX_SPREADSHEET_ROWS} rows)")
                        break
                    cells = [str(cell) if cell is not None else "" for cell in row]
                    rows.append(" | ".join(cells))
                if rows:
                    sheets.append(f"--- Sheet: {sheet_name} ---\n" + "\n".join(rows))
            wb.close()
            return "\n\n".join(sheets)[:MAX_TEXT_CHARS]
    except Exception as e:
        logger.warning("spreadsheet_extraction_failed", error=str(e))
        return f"[Spreadsheet extraction failed: {e}]"


async def process_file(
    file_info: Dict[str, Any],
    bot_token: str,
) -> Optional[FileAttachment]:
    """Process a Slack file attachment into content blocks for Claude.

    Args:
        file_info: Slack file object from event payload
        bot_token: Slack bot token for downloading

    Returns:
        FileAttachment with content blocks, or None if unsupported
    """
    filename = file_info.get("name", "unknown")
    mimetype = file_info.get("mimetype", "")
    file_url = file_info.get("url_private", "")
    size = file_info.get("size", 0)

    if not file_url:
        logger.warning("file_no_url", filename=filename)
        return None

    file_type = classify_file(filename, mimetype)
    logger.info(
        "processing_file", filename=filename, file_type=file_type, mimetype=mimetype, size=size
    )

    try:
        data, content_type = await download_slack_file(file_url, bot_token)
    except Exception as e:
        logger.error("file_download_failed", filename=filename, error=str(e))
        return None

    content_blocks: List[Dict[str, Any]] = []
    summary = ""

    if file_type == "image":
        # Send as base64 image block for Claude vision
        media_type = mimetype or content_type.split(";")[0]
        if media_type not in IMAGE_MIMES:
            media_type = "image/png"  # safe fallback

        # Compress/resize if needed to stay under Anthropic's 5MB base64 limit
        data, media_type = _compress_image_for_api(data, media_type)

        b64_data = base64.b64encode(data).decode("utf-8")
        content_blocks.append(
            {
                "type": "image",
                "source": {
                    "type": "base64",
                    "media_type": media_type,
                    "data": b64_data,
                },
            }
        )
        summary = f"Image: {filename} ({len(data)} bytes)"

    elif file_type == "pdf":
        text = extract_text_from_pdf(data)
        content_blocks.append(
            {
                "type": "text",
                "text": f"[Content of uploaded PDF: {filename}]\n\n{text}",
            }
        )
        summary = f"PDF: {filename} ({len(text)} chars extracted)"

    elif file_type == "spreadsheet":
        text = extract_text_from_spreadsheet(data, filename)
        content_blocks.append(
            {
                "type": "text",
                "text": f"[Content of uploaded spreadsheet: {filename}]\n\n{text}",
            }
        )
        summary = f"Spreadsheet: {filename} ({len(text)} chars extracted)"

    elif file_type == "document":
        text = extract_text_from_docx(data)
        content_blocks.append(
            {
                "type": "text",
                "text": f"[Content of uploaded document: {filename}]\n\n{text}",
            }
        )
        summary = f"Document: {filename} ({len(text)} chars extracted)"

    elif file_type == "presentation":
        text = extract_text_from_pptx(data)
        content_blocks.append(
            {
                "type": "text",
                "text": f"[Content of uploaded presentation: {filename}]\n\n{text}",
            }
        )
        summary = f"Presentation: {filename} ({len(text)} chars extracted)"

    elif file_type == "text":
        text = data.decode("utf-8", errors="replace")[:MAX_TEXT_CHARS]
        content_blocks.append(
            {
                "type": "text",
                "text": f"[Content of uploaded file: {filename}]\n\n{text}",
            }
        )
        summary = f"Text: {filename} ({len(text)} chars)"

    else:
        # Unknown file type — describe it
        content_blocks.append(
            {
                "type": "text",
                "text": f"[User uploaded file: {filename} ({mimetype}, {size} bytes). Unable to extract content from this file type.]",
            }
        )
        summary = f"Unknown: {filename} ({mimetype})"

    return FileAttachment(
        filename=filename,
        file_type=file_type,
        mime_type=mimetype,
        content_blocks=content_blocks,
        summary=summary,
        size_bytes=size,
    )


async def process_slack_files(
    files: List[Dict[str, Any]],
    bot_token: str,
) -> List[FileAttachment]:
    """Process multiple Slack file attachments.

    Args:
        files: List of Slack file objects from event payload
        bot_token: Slack bot token

    Returns:
        List of processed FileAttachments
    """
    attachments: List[FileAttachment] = []
    for file_info in files:
        attachment = await process_file(file_info, bot_token)
        if attachment:
            attachments.append(attachment)
    return attachments
